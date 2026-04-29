"""PowerPoint slide deck generation for the PA karst CATModel."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.config import OUT_FIGURES, OUT_PRESENTATION, OUT_TABLES, PROJECT_ROOT
from src.figures import DARK, SAGE, SAND, TERRACOTTA

WIDE_WIDTH = 13.333
WIDE_HEIGHT = 7.5
COURSE_FOOTER = "Lehigh University CAT402/CEE332 - PA Karst Hazard CATModel"
DASHBOARD_ASSETS = PROJECT_ROOT / "dashboard" / "assets"


def add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    line: str | None = None,
) -> Any:
    """Add a rectangle to a slide using inch coordinates and hex colors."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)
    return shape


def add_text(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 18,
    bold: bool = False,
    color: str = DARK,
    align: str = "left",
    anchor: str = "top",
    font: str = "Aptos",
    italic: bool = False,
) -> Any:
    """Add a single styled text box to a slide."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = _anchor(anchor)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = _align(align)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return box


def add_multi_text(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, int, bool, str, bool, int]],
    align: str = "left",
    anchor: str = "top",
    font: str = "Aptos",
) -> Any:
    """Add a multi-paragraph styled text box to a slide."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = _anchor(anchor)

    for idx, (text, size, bold, color, italic, space_after) in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.alignment = _align(align)
        paragraph.space_after = Pt(space_after)
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def build_presentation(out_path: Path) -> Path:
    """Build and save the final 15-slide project presentation."""
    prs = Presentation()
    prs.slide_width = Inches(WIDE_WIDTH)
    prs.slide_height = Inches(WIDE_HEIGHT)
    blank = prs.slide_layouts[6]

    _slide_cover(prs, blank)
    _slide_motivation(prs, blank, 2)
    _slide_study_area(prs, blank, 3)
    _slide_literature(prs, blank, 4)
    _slide_data_sources(prs, blank, 5)
    _slide_method1(prs, blank, 6)
    _slide_method2(prs, blank, 7)
    _slide_workflow(prs, blank, 8)
    _slide_result1(prs, blank, 9)
    _slide_result2(prs, blank, 10)
    _slide_result3(prs, blank, 11)
    _slide_result4(prs, blank, 12)
    _slide_gis_evidence(prs, blank, 13)
    _slide_limitations(prs, blank, 14)
    _slide_recommendations(prs, blank, 15)

    out_path = OUT_PRESENTATION / "PA_Karst_Final_Presentation.pptx" if out_path is None else out_path
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def _slide_cover(prs: Presentation, blank: Any) -> None:
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, WIDE_WIDTH, WIDE_HEIGHT, DARK)
    add_rect(slide, 0, 6.9, WIDE_WIDTH, 0.6, TERRACOTTA)
    add_multi_text(
        slide,
        0.8,
        1.05,
        9.6,
        2.5,
        [
            ("PA Karst", 38, True, SAND, False, 2),
            ("Hazard", 38, True, SAND, False, 2),
            ("CATModel", 38, True, SAND, False, 8),
        ],
        align="left",
    )
    add_text(slide, 0.85, 4.25, 5.8, 0.35, "Essential-facility exposure in southeastern Pennsylvania", 18, False, SAND)
    add_text(slide, 0.85, 5.0, 5.8, 0.35, "Lehigh University CAT402/CEE332 Final Project", 15, False, SAGE)
    add_text(slide, 0.85, 5.43, 5.8, 0.35, "Prepared by: N. Chatterjee", 15, False, SAGE)
    _notes(
        slide,
        "This presentation is the final package for the PA Karst Hazard CATModel. "
        "The project asks how mapped karst susceptibility overlaps essential facilities in five southeastern Pennsylvania counties. "
        "I frame the result as a screening model, not a parcel-level sinkhole prediction. "
        "The key deliverables are a reproducible Python pipeline, result tables, figures, this deck, and a Vercel project dashboard.",
    )


def _slide_motivation(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Motivation", page)
    add_text(slide, 0.8, 1.15, 6.0, 0.6, "Karst is a quiet hazard until infrastructure sits on top of it.", 26, True, TERRACOTTA)
    add_multi_text(
        slide,
        0.9,
        2.05,
        5.8,
        3.7,
        [
            ("Sinkholes and subsidence are spatially localized, often hydrologically triggered, and hard to convert into event probabilities at classroom-project scale.", 18, False, DARK, False, 12),
            ("A useful final project therefore needs a defensible screening model: map susceptibility, intersect essential facilities, and identify where follow-up engineering review would be most valuable.", 18, False, DARK, False, 12),
            ("The CATModel contribution is the explicit hazard-exposure-criticality framing and sensitivity testing.", 18, False, DARK, False, 0),
        ],
    )
    add_rect(slide, 7.4, 1.25, 4.7, 4.85, SAND, TERRACOTTA)
    add_text(slide, 7.8, 1.75, 3.9, 0.5, "Decision question", 22, True, DARK, "center")
    add_text(
        slide,
        7.85,
        2.65,
        3.8,
        2.3,
        "Where do mapped karst susceptibility and essential-facility exposure overlap strongly enough to justify deeper review?",
        21,
        True,
        TERRACOTTA,
        "center",
        "middle",
    )
    _notes(
        slide,
        "I start with the motivation because karst is not like a hurricane track or an earthquake rupture. "
        "The timing of a sinkhole at a specific parcel is difficult to defend from public map data alone. "
        "So the project deliberately becomes a screening model. "
        "The useful question is where mapped susceptibility and important facilities overlap enough to prioritize a closer look.",
    )


def _slide_study_area(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Study Area", page)
    stats = [
        ("5", "counties"),
        ("30,623", "study-area karst features"),
        ("1,717", "real facilities"),
        ("455", "facilities within 1 km"),
    ]
    for idx, (value, label) in enumerate(stats):
        x = 0.7 + idx * 3.12
        add_rect(slide, x, 1.15, 2.7, 1.18, SAND, TERRACOTTA)
        add_text(slide, x + 0.15, 1.33, 2.4, 0.4, value, 27, True, TERRACOTTA, "center")
        add_text(slide, x + 0.18, 1.82, 2.32, 0.32, label, 12, False, DARK, "center")
    summary = _read_table(OUT_TABLES / "table3_county_summary.csv")
    add_text(slide, 0.85, 2.8, 5.8, 0.42, "County exposure profile", 22, True, TERRACOTTA)
    add_text(slide, 0.95, 3.35, 5.5, 2.7, _county_summary_text(summary), 15, False, DARK)
    add_text(slide, 7.15, 2.85, 4.9, 0.42, "Spatial pattern", 22, True, TERRACOTTA)
    add_multi_text(
        slide,
        7.25,
        3.4,
        4.65,
        2.5,
        [
            ("Lehigh has the highest mapped karst density at 11.24 features/km2.", 16, False, DARK, False, 10),
            ("Berks has the largest mapped feature count at 16,294.", 16, False, DARK, False, 10),
            ("Bucks and Montgomery have many facilities but lower mapped karst densities.", 16, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "The study area is five counties in southeastern Pennsylvania. "
        "The regenerated pipeline finds 30,623 mapped karst features inside the county polygons. "
        "It links those hazards to 1,717 real facility points from DOH and USGS data. "
        "The key result is that 455 facilities fall within one kilometer of a mapped karst point.",
    )


def _slide_literature(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Literature Context", page)
    rows = [
        ("IPCC 2014", "Risk as interaction of hazard, exposure, and vulnerability."),
        ("Wood et al. 2023, USGS", "National-scale sinkhole susceptibility and exposure framing."),
        ("Reese & Kochanov, PAGS", "Pennsylvania karst and sinkhole engineering-geology context."),
        ("Galve et al. / Engineering Geology", "Sinkhole susceptibility and risk mapping precedent."),
        ("Qiu, Wu & Chen 2020", "GIS and imagery-derived sinkhole susceptibility attributes."),
    ]
    for idx, (source, role) in enumerate(rows):
        y = 1.25 + idx * 0.9
        add_rect(slide, 0.8, y, 2.7, 0.62, SAND, TERRACOTTA)
        add_text(slide, 0.95, y + 0.16, 2.4, 0.25, source, 12, True, DARK, "center")
        add_text(slide, 3.75, y + 0.1, 8.3, 0.38, role, 15, False, DARK)
    add_text(slide, 0.85, 6.2, 11.2, 0.35, "Takeaway: this project is a transparent screening model, not a calibrated event-probability model.", 16, True, TERRACOTTA, "center")
    _notes(
        slide,
        "The literature gives the project its boundaries. "
        "IPCC provides the risk framing, while Wood and colleagues show how sinkhole susceptibility can be handled at broad scale. "
        "Pennsylvania Geological Survey materials keep the model grounded in local karst behavior. "
        "The susceptibility mapping literature supports the idea that a map-based screen is useful but should not be oversold as prediction.",
    )


def _slide_data_sources(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Data Sources", page)
    rows = [
        ("DCNR / PASDA karst DBF", "144,245 statewide points; 30,623 in study counties"),
        ("PA DOH hospitals", "42 hospital points in the five-county study area"),
        ("USGS Structures", "1,675 school, fire/EMS, police, and correctional points"),
        ("PennDOT counties", "5 study county polygons and area values"),
        ("PennDOT municipalities", "287 municipal polygons for KHI aggregation"),
    ]
    for idx, (name, note) in enumerate(rows):
        y = 1.15 + idx * 0.86
        add_text(slide, 0.9, y, 3.8, 0.32, name, 15, True, TERRACOTTA)
        add_text(slide, 4.85, y, 7.0, 0.32, note, 15, False, DARK)
        add_rect(slide, 0.85, y + 0.46, 11.35, 0.02, SAND)
    add_multi_text(
        slide,
        0.9,
        5.75,
        11.1,
        0.9,
        [
            ("Implementation note: shapefiles are parsed with pure Python readers, Web Mercator boundaries are converted to lon/lat, and the pipeline falls back to fixed-seed synthetic facilities only when real facility files are absent.", 14, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "This slide is the reproducibility promise. "
        "The model uses public point inventories and PennDOT polygons rather than hand-entered locations. "
        "The karst DBF is parsed directly, and county and municipality polygons are used for assignment. "
        "If a clean clone only has the karst DBF, the code still runs but labels facility outputs as screening-grade synthetic.",
    )


def _slide_intro(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Introduction & Study Area", page)
    stats = [
        ("5", "counties"),
        ("30,623", "karst features"),
        ("1,717", "essential facilities"),
        ("287", "municipalities"),
    ]
    for idx, (value, label) in enumerate(stats):
        x = 0.8 + idx * 3.05
        add_rect(slide, x, 1.25, 2.55, 1.2, SAND, TERRACOTTA)
        add_text(slide, x + 0.18, 1.42, 2.2, 0.4, value, 28, True, TERRACOTTA, "center")
        add_text(slide, x + 0.18, 1.9, 2.2, 0.3, label, 13, False, DARK, "center")
    add_multi_text(
        slide,
        0.9,
        3.0,
        11.3,
        2.7,
        [
            ("Karst terrain creates localized ground-failure susceptibility that can matter disproportionately for hospitals, schools, emergency services, and police facilities.", 20, False, DARK, False, 12),
            ("This project builds a transparent CATModel-style screening workflow using mapped DCNR karst features, PA DOH hospitals, USGS essential structures, proximity buffers, and municipality-scale KHI rankings.", 18, False, DARK, False, 0),
        ],
    )


def _slide_method1(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Methodology I", page)
    add_text(slide, 0.85, 1.15, 5.6, 0.45, "Hazard reframed as susceptibility proxy", 24, True, TERRACOTTA)
    add_multi_text(
        slide,
        0.95,
        1.9,
        5.4,
        4.25,
        [
            ("Mapped karst features are treated as evidence of susceptibility, not as deterministic sinkhole events.", 18, False, DARK, False, 12),
            ("The model keeps installation simple by using pure-Python DBF and shapefile readers instead of GeoPandas.", 18, False, DARK, False, 12),
            ("County and municipality assignment uses real PennDOT polygons when the boundary archives are available.", 18, False, DARK, False, 12),
            ("This preserves reproducibility while still supporting CATModel-style exposure and sensitivity analysis.", 18, False, DARK, False, 0),
        ],
    )
    add_rect(slide, 7.0, 1.35, 4.8, 4.4, SAND, TERRACOTTA)
    add_text(slide, 7.35, 1.75, 4.1, 0.55, "Susceptibility layer", 23, True, DARK, "center")
    add_text(slide, 7.35, 2.55, 4.1, 1.75, "DCNR mapped karst points\n+\nfacility proximity buffers", 24, True, TERRACOTTA, "center", "middle")
    add_text(slide, 7.35, 4.85, 4.1, 0.55, "screening, not prediction", 16, False, DARK, "center")
    _notes(
        slide,
        "The first methodology point is language discipline. "
        "A mapped karst point is evidence of susceptibility, not a timed event. "
        "The code uses real polygons when present, but it does not claim to model subsurface mechanics. "
        "This framing lets the results be useful without overstating what public data can prove.",
    )


def _slide_method2(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Methodology II", page)
    add_text(slide, 0.85, 1.2, 5.9, 0.5, "Karst Hazard Index", 25, True, TERRACOTTA)
    add_text(slide, 0.95, 2.0, 5.7, 0.9, "KHI = 100 * (wH * H + wE * E + wC * C)", 24, True, DARK, "center")
    add_multi_text(
        slide,
        0.95,
        3.25,
        5.7,
        2.2,
        [
            ("H: normalized municipal karst density", 17, False, DARK, False, 8),
            ("E: normalized exposed-facility count within 500 m", 17, False, DARK, False, 8),
            ("C: normalized criticality-weighted exposure", 17, False, DARK, False, 0),
        ],
    )
    add_rect(slide, 7.25, 1.35, 4.8, 4.4, SAND, TERRACOTTA)
    add_multi_text(
        slide,
        7.6,
        1.8,
        4.1,
        3.5,
        [
            ("Weighting schemes", 21, True, DARK, False, 16),
            ("Base: 0.40 / 0.40 / 0.20", 16, False, DARK, False, 8),
            ("Hazard-led: 0.60 / 0.25 / 0.15", 16, False, DARK, False, 8),
            ("Exposure-led: 0.25 / 0.60 / 0.15", 16, False, DARK, False, 8),
            ("Critical-facility: 0.30 / 0.35 / 0.35", 16, False, DARK, False, 18),
            ("Framing follows IPCC 2014, Wood et al. USGS, and Reese & Kochanov PAGS.", 13, False, DARK, True, 0),
        ],
    )
    _notes(
        slide,
        "The KHI is deliberately simple and auditable. "
        "Hazard is normalized municipal karst density, exposure is the count of facilities within 500 meters, and criticality weights facility function and proximity. "
        "I test four schemes so the ranking is not dependent on a single subjective weighting choice. "
        "The Spearman correlations later show whether the priority list is stable.",
    )


def _slide_workflow(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Workflow", page)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig7_flowchart.png", 0.75, 1.15, 11.75, 3.95)
    add_multi_text(
        slide,
        0.95,
        5.35,
        11.75,
        0.95,
        [
            ("Every step writes an inspectable artifact: processed CSVs, tables, figures, a PowerPoint deck, and a deployable dashboard. The pipeline can be run end to end with python scripts/run_all.py after raw data are placed in data/raw/.", 15, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "The workflow is intentionally linear and reproducible. "
        "First the karst DBF is parsed and filtered, then facilities are prepared, then exposure distances and buffers are computed. "
        "The municipality KHI and sensitivity analysis come next, followed by figures, tables, the deck, and the dashboard. "
        "This matters because every number shown later can be traced back to a generated table.",
    )


def _slide_result1(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Result I: Karst Density", page)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig1_karst_density.png", 0.65, 1.15, 7.3, 5.6)
    add_multi_text(
        slide,
        8.25,
        1.35,
        4.1,
        4.9,
        [
            ("Key findings", 22, True, TERRACOTTA, False, 14),
            ("Karst features cluster unevenly across the study area.", 17, False, DARK, False, 10),
            ("The density surface is a screening layer for exposure, not a probability surface.", 17, False, DARK, False, 10),
            ("County bounding boxes provide a simple reproducible geography for the final project.", 17, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "This figure shows that the hazard proxy is highly clustered. "
        "The visual concentration is strongest in Lehigh and Berks, which also dominate later exposure results. "
        "I describe this as density of mapped features, not probability of failure. "
        "The point is to identify where susceptibility should receive more attention.",
    )


def _slide_result2(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Result II: Facility Exposure", page)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig2_facility_exposure.png", 0.65, 1.1, 7.5, 5.65)
    summary = _read_table(OUT_TABLES / "table3_county_summary.csv")
    text = _county_summary_text(summary)
    add_multi_text(
        slide,
        8.45,
        1.35,
        3.85,
        4.8,
        [
            ("Exposure summary", 22, True, TERRACOTTA, False, 14),
            (text, 16, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "This slide connects the hazard layer to actual facilities. "
        "The regenerated exposure model finds 455 facilities within one kilometer of mapped karst. "
        "Lehigh has the highest exposure share, while Berks has a large absolute count. "
        "The map is useful for screening where asset QA and site review should come first.",
    )


def _slide_result3(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Result III: Facility Risk", page)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig5_top10_facilities.png", 0.55, 1.15, 6.25, 5.25)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig4_county_summary.png", 6.95, 1.15, 5.8, 5.25)
    _notes(
        slide,
        "The top facility chart ranks individual records by local density, proximity, and criticality. "
        "Schools and fire facilities can appear high because they are numerous and sometimes sit directly on dense mapped clusters. "
        "The county summary adds scale so we do not overinterpret one facility at a time. "
        "Together these two views show both asset-level and county-level priorities.",
    )


def _slide_result4(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Result IV: KHI Sensitivity", page)
    _add_picture_or_placeholder(slide, OUT_FIGURES / "fig6_sensitivity.png", 0.55, 1.15, 7.4, 4.9)
    top5 = _top5_munis()
    add_multi_text(
        slide,
        8.25,
        1.25,
        4.25,
        5.35,
        [
            ("Top 5 Base KHI", 21, True, TERRACOTTA, False, 10),
            (top5, 14, False, DARK, False, 14),
            ("Uncertainty", 18, True, TERRACOTTA, False, 8),
            ("Rank sensitivity indicates whether priority cells are robust to alternative hazard-led or exposure-led assumptions.", 14, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "The KHI sensitivity slide tests whether the priority ranking is fragile. "
        "Allentown remains the top base-scheme municipality because it combines many exposed facilities with meaningful karst density. "
        "Hanover and Lower Macungie move depending on how much weight is placed on hazard density. "
        "This is why rank stability is part of the final interpretation, not an afterthought.",
    )


def _slide_gis_evidence(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "GIS Evidence", page)
    _add_picture_or_placeholder(slide, DASHBOARD_ASSETS / "arcgis-clip-workflow.png", 0.62, 1.12, 6.35, 3.15)
    _add_picture_or_placeholder(slide, DASHBOARD_ASSETS / "county-exposure-layout.png", 7.22, 1.12, 2.45, 3.15)
    _add_picture_or_placeholder(slide, DASHBOARD_ASSETS / "karst-density-layout.png", 9.88, 1.12, 2.45, 3.15)
    add_multi_text(
        slide,
        0.72,
        4.62,
        11.75,
        1.45,
        [
            ("Production record", 20, True, TERRACOTTA, False, 8),
            ("The ArcGIS Pro clip screenshot and exported layouts document the preprocessing chain used before final Python/Vercel packaging. They are included as evidence artifacts rather than hidden intermediate work.", 14, False, DARK, False, 0),
        ],
    )
    _notes(
        slide,
        "This slide is included because the project also has a GIS production history. "
        "The screenshot shows the clip workflow in ArcGIS Pro, while the two layouts show the mapped hazard proxy and density communication layers. "
        "These artifacts help a professor see that the project was not only a web dashboard. "
        "They also connect the Python outputs to the GIS workflow used during development.",
    )


def _slide_limitations(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Limitations & Uncertainty", page)
    items = [
        ("Hazard proxy", "Mapped karst is evidence of susceptibility, not a timed failure probability."),
        ("Geometry assignment", "Polygon assignment is real, but boundary/source coordinate error still exists."),
        ("Facilities", "Public inventories require facility-by-facility QA before operational use."),
        ("Vulnerability", "No fragility curves, downtime, service areas, or repair-cost model are included."),
        ("Temporal risk", "No rainfall, groundwater, construction, or climate trigger probability is modeled."),
    ]
    for idx, (title, body) in enumerate(items):
        x = 0.8 + (idx % 3) * 4.15
        y = 1.25 + (idx // 3) * 2.25
        add_rect(slide, x, y, 3.55, 1.55, SAND, TERRACOTTA)
        add_text(slide, x + 0.22, y + 0.22, 3.1, 0.3, title, 15, True, TERRACOTTA)
        add_text(slide, x + 0.22, y + 0.62, 3.05, 0.7, body, 12, False, DARK)
    add_text(slide, 0.9, 6.35, 11.3, 0.35, "Interpretation rule: results identify screening priorities, not evidence of current facility damage.", 15, True, DARK, "center")
    _notes(
        slide,
        "This is the slide that protects the project from overclaiming. "
        "The model is strong as a transparent screen, but it does not contain site geotechnical investigation or temporal probability. "
        "It also does not estimate structural damage or service downtime. "
        "Those limitations become the future-work roadmap rather than a weakness hidden from the audience.",
    )


def _slide_recommendations(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Recommendations & Conclusion", page)
    columns = [
        ("Immediate QA", "Verify top hospital and high-risk facility coordinates against source records, aerial imagery, and local asset owners."),
        ("Engineering upgrade", "Add carbonate geology, depth-to-bedrock, drainage, slope, soils, and stormwater covariates."),
        ("CATModel upgrade", "Replace relative loss with fragility, downtime, service disruption, and dollar-denominated repair/loss estimates."),
    ]
    for idx, (title, body) in enumerate(columns):
        x = 0.75 + idx * 4.18
        add_rect(slide, x, 1.25, 3.55, 3.25, SAND, TERRACOTTA)
        add_text(slide, x + 0.25, 1.65, 3.05, 0.4, title, 17, True, TERRACOTTA, "center")
        add_text(slide, x + 0.32, 2.25, 2.9, 1.55, body, 13, False, DARK, "center", "middle")
    add_text(
        slide,
        1.0,
        5.2,
        11.3,
        0.8,
        "Bottom line: the strongest screening priorities are where dense mapped karst, essential facilities, and rank-stable KHI results converge.",
        23,
        True,
        TERRACOTTA,
        "center",
        "middle",
    )
    _notes(
        slide,
        "My closing recommendation is to treat the model as a triage tool. "
        "The first practical step is QA on the highest-ranked facilities and municipalities. "
        "The next engineering step is to add geotechnical and hydrologic covariates. "
        "The next CATModel step is to turn exposure into damage, downtime, and loss.",
    )


def _slide_conclusion(prs: Presentation, blank: Any, page: int) -> None:
    slide = _content_slide(prs, blank, "Conclusion", page)
    columns = [
        ("CATModel takeaways", "A lightweight hazard, exposure, and criticality workflow can identify where karst screening deserves closer engineering attention."),
        ("Limitations", "Mapped karst is a susceptibility proxy, public facility coordinates require QA, and no fragility curves or service disruption model are included."),
        ("Future work", "Add facility QA, service-area geometry, geotechnical covariates, and calibrated event likelihood or damage functions."),
    ]
    for idx, (title, body) in enumerate(columns):
        x = 0.75 + idx * 4.18
        add_rect(slide, x, 1.35, 3.55, 4.85, SAND, TERRACOTTA)
        add_text(slide, x + 0.25, 1.75, 3.05, 0.55, title, 19, True, TERRACOTTA, "center")
        add_text(slide, x + 0.35, 2.65, 2.85, 2.8, body, 15, False, DARK, "center", "middle")


def _content_slide(prs: Presentation, blank: Any, title: str, page: int) -> Any:
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, WIDE_WIDTH, WIDE_HEIGHT, "FFFFFF")
    add_rect(slide, 0, 0, WIDE_WIDTH, 0.22, TERRACOTTA)
    add_text(slide, 0.65, 0.42, 9.5, 0.45, title, 24, True, DARK)
    _footer(slide, page)
    return slide


def _footer(slide: Any, page: int) -> None:
    add_rect(slide, 0, 7.12, WIDE_WIDTH, 0.38, SAND)
    add_text(slide, 0.55, 7.18, 9.5, 0.2, COURSE_FOOTER, 8, False, DARK)
    add_text(slide, 12.2, 7.18, 0.6, 0.2, str(page), 8, False, DARK, "right")


def _add_picture_or_placeholder(
    slide: Any, path: Path, x: float, y: float, w: float, h: float
) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_rect(slide, x, y, w, h, SAND, TERRACOTTA)
        add_text(slide, x + 0.25, y + h / 2 - 0.2, w - 0.5, 0.4, path.name, 15, True, DARK, "center")


def _read_table(path: Path) -> pd.DataFrame:
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame()


def _county_summary_text(summary: pd.DataFrame) -> str:
    if summary.empty:
        return "Run the pipeline to populate county exposure summaries."
    parts = []
    for row in summary.itertuples(index=False):
        parts.append(
            f"{row.county}: {int(row.exposed_within_1000m)} of {int(row.total)} within 1 km"
        )
    return "\n".join(parts)


def _top5_munis() -> str:
    table = _read_table(OUT_TABLES / "table4_top15_munis_KHI.csv")
    if table.empty:
        return "Run the pipeline to populate KHI rankings."
    lines = []
    for idx, row in enumerate(table.head(5).itertuples(index=False), start=1):
        lines.append(f"{idx}. {row.label}: {row.KHI_Base:.1f}")
    return "\n".join(lines)


def _align(value: str) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(value, PP_ALIGN.LEFT)


def _anchor(value: str) -> MSO_ANCHOR:
    return {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(value, MSO_ANCHOR.TOP)


def _rgb(hex_color: str) -> RGBColor:
    color = hex_color.lstrip("#")
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _notes(slide: Any, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text
